from __future__ import annotations

import argparse
import datetime as dt
from io import BytesIO
import re
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

try:
    from PIL import Image, ImageDraw  # type: ignore
except Exception:  # pragma: no cover
    Image = None  # type: ignore
    ImageDraw = None  # type: ignore


@dataclass
class SubSection:
    title: str
    items: list[tuple[int, str]] = field(default_factory=list)  # (indent_level, text)


@dataclass
class Section:
    title: str
    subsections: list[SubSection] = field(default_factory=list)


HEADING_RE = re.compile(r"^(?P<level>#{1,6})\s+(?P<title>.+?)\s*$")
BULLET_RE = re.compile(r"^(?P<indent>\s*)-\s+(?P<text>.+?)\s*$")


def parse_requirements_md(md_text: str) -> tuple[str | None, list[Section]]:
    doc_title: str | None = None
    sections: list[Section] = []

    current_section: Section | None = None
    current_sub: SubSection | None = None

    def flush_sub() -> None:
        nonlocal current_sub, current_section
        if current_sub is None:
            return
        if current_section is None:
            current_section = Section(title="其他")
            sections.append(current_section)
        current_section.subsections.append(current_sub)
        current_sub = None

    def flush_section() -> None:
        nonlocal current_section
        flush_sub()
        current_section = None

    for raw_line in md_text.splitlines():
        line = raw_line.rstrip("\n")
        if not line.strip():
            continue

        hm = HEADING_RE.match(line)
        if hm:
            level = len(hm.group("level"))
            title = hm.group("title").strip()
            if level == 1:
                doc_title = title
                flush_section()
            elif level == 2:
                flush_section()
                current_section = Section(title=title)
                sections.append(current_section)
            elif level == 3:
                flush_sub()
                current_sub = SubSection(title=title)
            else:
                # Treat deeper headings as part of current subsection.
                if current_sub is None:
                    current_sub = SubSection(title=title)
                else:
                    current_sub.items.append((0, title))
            continue

        bm = BULLET_RE.match(line)
        if bm:
            indent_spaces = len(bm.group("indent").replace("\t", "    "))
            level = max(0, indent_spaces // 2)
            text = bm.group("text").strip()
            if current_sub is None:
                current_sub = SubSection(title="要点")
            current_sub.items.append((level, text))
            continue

        # Fallback: treat as a bullet-like line in current subsection.
        if current_sub is None:
            current_sub = SubSection(title="说明")
        current_sub.items.append((0, line.strip()))

    flush_section()
    return doc_title, sections


def chunk_items(items: list[tuple[int, str]], max_lines: int) -> list[list[tuple[int, str]]]:
    if not items:
        return [[]]
    chunks: list[list[tuple[int, str]]] = []
    buf: list[tuple[int, str]] = []
    for it in items:
        buf.append(it)
        if len(buf) >= max_lines:
            chunks.append(buf)
            buf = []
    if buf:
        chunks.append(buf)
    return chunks


def _estimate_lines(text: str, *, chars_per_line: int) -> int:
    t = text.strip()
    if not t:
        return 0
    return max(1, (len(t) + chars_per_line - 1) // chars_per_line)


def _pack_subsections_into_slides(
    subsections: list[SubSection],
    *,
    max_lines: int,
    header_line_cost: int = 1,
    chars_per_line: int = 42,
) -> list[list[tuple[str, list[tuple[int, str]]]]]:
    """
    Return a list of slides, each slide is a list of (sub_title, items_chunk).
    """
    slides: list[list[tuple[str, list[tuple[int, str]]]]] = []
    current: list[tuple[str, list[tuple[int, str]]]] = []
    used_lines = 0

    def flush() -> None:
        nonlocal current, used_lines
        if current:
            slides.append(current)
            current = []
            used_lines = 0

    for sub in subsections:
        # Estimate per-item line usage with a rough wrap heuristic.
        item_line_costs = [max(1, _estimate_lines(t, chars_per_line=chars_per_line)) for _, t in sub.items]
        sub_total = header_line_cost + sum(item_line_costs)

        if sub_total <= max_lines:
            if used_lines + sub_total > max_lines:
                flush()
            current.append((sub.title, sub.items))
            used_lines += sub_total
            continue

        # Subsection too large: split across slides.
        remaining = list(sub.items)
        part = 1
        while remaining:
            if used_lines + header_line_cost + 1 > max_lines:
                flush()

            available = max_lines - used_lines - header_line_cost
            if available <= 0:
                flush()
                available = max_lines - header_line_cost

            chunk: list[tuple[int, str]] = []
            used = 0
            while remaining and used < available:
                lvl, txt = remaining[0]
                cost = max(1, _estimate_lines(txt, chars_per_line=chars_per_line))
                if chunk and used + cost > available:
                    break
                chunk.append((lvl, txt))
                used += cost
                remaining.pop(0)

            # Keep slide headers clean: do not append “续n” by default.
            title = sub.title
            current.append((title, chunk))
            used_lines += header_line_cost + used
            part += 1
            if remaining:
                flush()

    flush()
    return slides


def _set_run_font(run, *, name: str, size_pt: int, bold: bool | None = None, color: RGBColor | None = None) -> None:
    run.font.name = name
    # Ensure East Asian font is set for Chinese text. python-pptx doesn't expose this directly.
    try:
        r_pr = run._r.get_or_add_rPr()
        r_fonts = r_pr.get_or_add_rFonts()
        r_fonts.set(qn("w:eastAsia"), name)
    except Exception:
        pass
    run.font.size = Pt(size_pt)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _set_paragraph_font(paragraph, *, name: str, size_pt: int, bold: bool | None = None, color: RGBColor | None = None) -> None:
    for run in paragraph.runs:
        _set_run_font(run, name=name, size_pt=size_pt, bold=bold, color=color)
    if not paragraph.runs and paragraph.text:
        run = paragraph.add_run()
        run.text = paragraph.text
        paragraph.text = ""
        _set_run_font(run, name=name, size_pt=size_pt, bold=bold, color=color)


def _add_footer(slide, text: str, *, font_name: str) -> None:
    left = Inches(0.3)
    top = Inches(7.05)
    width = Inches(12.7)
    height = Inches(0.35)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = 2  # right
    run = p.add_run()
    run.text = text
    _set_run_font(run, name=font_name, size_pt=10, color=RGBColor(120, 120, 120))


def _mix_rgb(fg: tuple[int, int, int], bg: tuple[int, int, int], alpha_pct: int) -> tuple[int, int, int]:
    a = max(0, min(100, alpha_pct))
    return tuple(int((f * a + b * (100 - a)) / 100) for f, b in zip(fg, bg))


def _save_png(path: Path, img) -> None:
    buf = BytesIO()
    img.save(buf, format="PNG")
    path.write_bytes(buf.getvalue())


def _draw_card(draw, *, x0: int, y0: int, x1: int, y1: int, fill: tuple[int, int, int], outline: tuple[int, int, int]) -> None:
    try:
        draw.rounded_rectangle([x0, y0, x1, y1], radius=26, fill=fill, outline=outline, width=3)
    except Exception:
        draw.rectangle([x0, y0, x1, y1], fill=fill, outline=outline, width=3)


def _write_theme_backgrounds(asset_dir: Path, theme: str) -> tuple[Path, Path]:
    """
    Write (or overwrite) a theme's title + normal backgrounds.
    Returns (title_bg_path, normal_bg_path).
    """
    theme_dir = asset_dir / "themes" / theme
    theme_dir.mkdir(parents=True, exist_ok=True)
    title_bg = theme_dir / "bg_title.png"
    normal_bg = theme_dir / "bg_normal.png"

    if Image is None or ImageDraw is None:
        return title_bg, normal_bg

    w, h = 1920, 1080

    # Theme palette
    palettes: dict[str, dict[str, tuple[int, int, int]]] = {
        "light-blue": {
            "accent": (31, 78, 121),
            "accent_light": (234, 242, 251),
            "bg": (247, 250, 255),
            "border": (220, 230, 240),
        },
        "light-gray": {
            "accent": (55, 55, 55),
            "accent_light": (244, 244, 244),
            "bg": (250, 250, 250),
            "border": (230, 230, 230),
        },
        "minimal": {
            "accent": (55, 55, 55),
            "accent_light": (250, 250, 250),
            "bg": (255, 255, 255),
            "border": (235, 235, 235),
        },
        "light-green": {
            "accent": (27, 94, 58),
            "accent_light": (232, 246, 238),
            "bg": (246, 252, 248),
            "border": (220, 236, 228),
        },
        "light-teal": {
            "accent": (0, 102, 102),
            "accent_light": (231, 248, 248),
            "bg": (246, 253, 253),
            "border": (215, 235, 235),
        },
        "light-purple": {
            "accent": (76, 46, 140),
            "accent_light": (242, 238, 252),
            "bg": (250, 249, 255),
            "border": (232, 228, 245),
        },
        "light-orange": {
            "accent": (170, 82, 0),
            "accent_light": (255, 242, 229),
            "bg": (255, 252, 248),
            "border": (245, 228, 210),
        },
        "sunset": {
            "accent": (30, 30, 30),
            "accent_light": (255, 238, 232),
            "bg": (255, 250, 248),
            "border": (244, 222, 215),
        },
    }
    pal = palettes.get(theme, palettes["light-blue"])
    accent = pal["accent"]
    accent_light = pal["accent_light"]
    bg = pal["bg"]
    border = pal["border"]

    # Normal slide background
    img = Image.new("RGB", (w, h), bg)
    draw = ImageDraw.Draw(img)
    top_h = 128
    bottom_h = 64

    # Light title bar so black titles remain readable.
    draw.rectangle([0, 0, w, top_h], fill=accent_light)
    top_line = pal["accent"] if theme not in {"light-gray", "minimal"} else (120, 120, 120)
    draw.rectangle([0, top_h - 6, w, top_h], fill=top_line)
    draw.rectangle([0, h - bottom_h, w, h], fill=(255, 255, 255) if theme != "minimal" else (255, 255, 255))
    draw.line([0, top_h, w, top_h], fill=border, width=3)
    draw.line([0, h - bottom_h, w, h - bottom_h], fill=border, width=3)

    if theme == "minimal":
        # Minimal: no card outline, only a subtle frame.
        draw.rectangle([90, top_h + 48, w - 90, h - bottom_h - 46], outline=border, width=2)
    else:
        _draw_card(draw, x0=90, y0=top_h + 48, x1=w - 90, y1=h - bottom_h - 46, fill=(255, 255, 255), outline=border)

    if theme == "sunset":
        # Gradient-ish bands
        draw.rectangle([0, top_h, w, top_h + 10], fill=(255, 180, 165))
        draw.rectangle([0, top_h + 10, w, top_h + 20], fill=(255, 214, 170))

    # Subtle circles pattern on the right
    ring_color = pal["accent"] if theme not in {"light-gray", "minimal"} else (120, 120, 120)
    for r, a in [(320, 10), (240, 14), (180, 18)]:
        mixed = _mix_rgb(ring_color, bg, a)
        cx, cy = int(w * 0.82), top_h + 210
        draw.ellipse([cx - r, cy - r, cx + r, cy + r], outline=mixed, width=6)

    _save_png(normal_bg, img)

    # Title slide background
    img2 = Image.new("RGB", (w, h), (255, 255, 255))
    draw2 = ImageDraw.Draw(img2)

    draw2.polygon([(0, 0), (w, 0), (w, int(h * 0.42)), (0, int(h * 0.18))], fill=_mix_rgb(accent_light, (255, 255, 255), 100))
    draw2.polygon([(0, 0), (int(w * 0.62), 0), (0, int(h * 0.34))], fill=_mix_rgb(accent_light, (255, 255, 255), 100))
    bar_color = ring_color if theme not in {"light-gray", "minimal"} else (120, 120, 120)
    draw2.rectangle([0, 0, int(w * 0.62), 8], fill=bar_color)
    draw2.rectangle([0, h - 8, w, h], fill=bar_color)
    for r, a in [(360, 6), (260, 8), (200, 10)]:
        mixed = _mix_rgb(bar_color, (255, 255, 255), a)
        cx, cy = int(w * 0.78), int(h * 0.55)
        draw2.ellipse([cx - r, cy - r, cx + r, cy + r], outline=mixed, width=6)

    _save_png(title_bg, img2)
    return title_bg, normal_bg


def ensure_builtin_themes(asset_dir: Path) -> list[str]:
    themes = ["light-blue", "light-gray", "minimal", "light-green", "light-teal", "light-purple", "light-orange", "sunset"]
    for t in themes:
        _write_theme_backgrounds(asset_dir, t)
    return themes


def _add_background_image(prs: Presentation, slide, bg_path: Path) -> None:
    if not bg_path.exists():
        return
    pic = slide.shapes.add_picture(str(bg_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    try:
        sp_tree = slide.shapes._spTree  # noqa: SLF001
        sp_tree.remove(pic._element)
        sp_tree.insert(2, pic._element)
    except Exception:
        pass


def _add_footer_lr(slide, *, left_text: str, right_text: str, font_name: str) -> None:
    top = Inches(7.05)
    height = Inches(0.35)

    left_box = slide.shapes.add_textbox(Inches(0.6), top, Inches(6.2), height)
    _force_horizontal_text(left_box)
    tf = left_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = left_text
    _set_run_font(run, name=font_name, size_pt=10, color=RGBColor(120, 120, 120))

    right_box = slide.shapes.add_textbox(Inches(7.0), top, Inches(5.9), height)
    _force_horizontal_text(right_box)
    tf2 = right_box.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = right_text
    _set_run_font(run2, name=font_name, size_pt=10, color=RGBColor(120, 120, 120))


def _force_horizontal_text(shape) -> None:
    """
    PowerPoint can store text direction as vertical in shape body properties.
    Force horizontal direction and reset rotation where possible.
    """
    try:
        shape.rotation = 0
    except Exception:
        pass
    try:
        tf = shape.text_frame
    except Exception:
        return
    try:
        # Some viewers are sensitive to text-direction attributes. Prefer removing them.
        # If the layout explicitly marks vertical text, override it to horizontal.
        vert = tf._bodyPr.get("vert")
        vertical_vals = {
            "eaVert",
            "mongolianVert",
            "vert",
            "vert270",
            "wordArtVert",
            "wordArtVertRtl",
        }
        if vert in vertical_vals:
            tf._bodyPr.set("vert", "horz")
        else:
            tf._bodyPr.attrib.pop("vert", None)
        tf._bodyPr.attrib.pop("upright", None)
        tf._bodyPr.attrib.pop("rot", None)
    except Exception:
        pass


def _hide_shape(shape) -> None:
    try:
        shape.left = Inches(-30)
        shape.top = Inches(-30)
    except Exception:
        pass


def _add_custom_title(
    slide,
    *,
    text: str,
    font_name: str,
    font_size_pt: int,
    color: RGBColor,
    left_in: float = 0.6,
    top_in: float = 0.12,
    width_in: float = 12.2,
    height_in: float = 0.7,
) -> None:
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    _force_horizontal_text(box)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    _set_run_font(run, name=font_name, size_pt=font_size_pt, bold=True, color=color)


def _shift_title_up(slide, *, font_size_pt: int, pct: float = 0.60, min_top_in: float = 0.12) -> None:
    """
    Shift the slide title placeholder upward.
    If we can't see the background image geometry, approximate by % of font height.
    """
    if slide.shapes.title is None:
        return
    try:
        shift_in = (font_size_pt * pct) / 72.0
        new_top = max(Inches(min_top_in), slide.shapes.title.top - Inches(shift_in))
        slide.shapes.title.top = new_top
    except Exception:
        pass


def build_ppt(
    sections: list[Section],
    *,
    title: str,
    subtitle: str,
    output_path: Path,
    bg_path: Path | None = None,
    bg_title_path: Path | None = None,
    theme: str | None = "light-blue",
) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    font_name = "Microsoft YaHei"
    accent = RGBColor(31, 78, 121)
    today = dt.date.today().isoformat()

    asset_dir = output_path.parent / ".ppt_assets"
    ensure_builtin_themes(asset_dir)

    normal_bg = bg_path
    title_bg = bg_title_path
    if normal_bg is None and title_bg is None:
        chosen = theme
        if not chosen:
            chosen = "light-blue"
        t_bg, n_bg = _write_theme_backgrounds(asset_dir, chosen)
        normal_bg = n_bg
        title_bg = t_bg
    else:
        # If only one provided, use it for both.
        if normal_bg is None and title_bg is not None:
            normal_bg = title_bg
        if title_bg is None and normal_bg is not None:
            title_bg = normal_bg

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _add_background_image(prs, slide, title_bg)
    if slide.shapes.title is not None:
        _hide_shape(slide.shapes.title)
    _add_custom_title(slide, text=title, font_name=font_name, font_size_pt=44, color=RGBColor(0, 0, 0), top_in=0.22)
    if len(slide.placeholders) > 1:
        subtitle_ph = slide.placeholders[1]
        subtitle_ph.text = subtitle
        _force_horizontal_text(subtitle_ph)
        spara = subtitle_ph.text_frame.paragraphs[0]
        spara.alignment = PP_ALIGN.LEFT
        _set_paragraph_font(spara, name=font_name, size_pt=18, color=RGBColor(80, 80, 80))

    # Agenda
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_background_image(prs, slide, normal_bg)
    if slide.shapes.title is not None:
        _hide_shape(slide.shapes.title)
    _add_custom_title(slide, text="目录", font_name=font_name, font_size_pt=30, color=RGBColor(0, 0, 0), top_in=0.16)
    body_ph = slide.shapes.placeholders[1]
    _force_horizontal_text(body_ph)
    body = body_ph.text_frame
    body.clear()
    body.word_wrap = True
    for i, sec in enumerate(sections, start=1):
        p = body.paragraphs[0] if i == 1 else body.add_paragraph()
        p.text = f"{i}. {sec.title}"
        p.level = 0
        _set_paragraph_font(p, name=font_name, size_pt=20, color=RGBColor(40, 40, 40))

    # Sections + subsections
    for sec in sections:
        packed = _pack_subsections_into_slides(sec.subsections, max_lines=18)
        for idx, blocks in enumerate(packed, start=1):
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only (we add our own body textbox)
            _add_background_image(prs, slide, normal_bg)
            if slide.shapes.title is not None:
                _hide_shape(slide.shapes.title)
            _add_custom_title(slide, text=sec.title, font_name=font_name, font_size_pt=28, color=RGBColor(0, 0, 0), top_in=0.16)

            body_left = Inches(0.8)
            base_top = 1.35
            base_height = 5.6
            body_width = Inches(11.8)

            # Balance sparse pages: if only a couple of short subsections, center content vertically a bit.
            def estimate_slide_lines() -> int:
                total = 0
                for sub_title, items in blocks:
                    total += 1  # header
                    for _, t in items:
                        total += max(1, _estimate_lines(t, chars_per_line=42))
                return total

            max_lines = 18
            used_lines = estimate_slide_lines()
            ratio = used_lines / max_lines if max_lines else 1.0
            shift = 0.0
            if len(blocks) <= 2 and ratio <= 0.55:
                # Up to ~0.7in shift-down when very sparse.
                shift = min(0.7, max(0.0, (0.55 - ratio) * 1.6))

            body_top = Inches(base_top + shift)
            body_height = Inches(max(4.2, base_height - shift))
            box = slide.shapes.add_textbox(body_left, body_top, body_width, body_height)
            _force_horizontal_text(box)
            tf = box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.margin_left = Inches(0.15)
            tf.margin_right = Inches(0.12)
            tf.margin_top = Inches(0.08)
            tf.margin_bottom = Inches(0.08)

            first_para = True
            for sub_title, items in blocks:
                p = tf.paragraphs[0] if first_para else tf.add_paragraph()
                first_para = False
                p.text = sub_title
                p.level = 0
                _set_paragraph_font(p, name=font_name, size_pt=20, bold=True, color=accent)
                try:
                    p.space_after = Pt(4)
                except Exception:
                    pass

                if not items:
                    p2 = tf.add_paragraph()
                    p2.text = "（无）"
                    p2.level = 1
                    _set_paragraph_font(p2, name=font_name, size_pt=16, color=RGBColor(80, 80, 80))
                else:
                    for lvl, text in items:
                        bp = tf.add_paragraph()
                        bp.text = text
                        bp.level = min(5, 1 + max(0, lvl))
                        _set_paragraph_font(bp, name=font_name, size_pt=16, color=RGBColor(40, 40, 40))
                        try:
                            bp.space_after = Pt(1)
                        except Exception:
                            pass

    output_path.parent.mkdir(parents=True, exist_ok=True)
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, start=1):
        left = today if i == 1 else f"{today} · {title}"
        right = "" if i == 1 else f"第 {i-1} / {max(1, total-1)} 页"
        _add_footer_lr(slide, left_text=left, right_text=right, font_name=font_name)
    prs.save(str(output_path))


def main() -> int:
    parser = argparse.ArgumentParser(description="Generate a PPTX from requirements.md")
    parser.add_argument("--input", default="requirements.md", help="Path to requirements.md")
    parser.add_argument("--output", default=str(Path("doc") / "智能派单需求说明.pptx"), help="Output .pptx path")
    parser.add_argument("--bg", default="", help="Background image path for normal slides (png/jpg).")
    parser.add_argument("--bg-title", default="", help="Background image path for title slide (png/jpg).")
    parser.add_argument("--theme", default="light-blue", help="Built-in theme name (see --list-themes).")
    parser.add_argument("--list-themes", action="store_true", help="Print built-in themes and write preview assets.")
    args = parser.parse_args()

    in_path = Path(args.input)
    out_path = Path(args.output)

    md_text = in_path.read_text(encoding="utf-8")
    doc_title, sections = parse_requirements_md(md_text)

    title = "智能派单需求说明"
    if doc_title and doc_title.strip():
        title = doc_title.strip()
    subtitle = f"来源：{in_path.name}（自动生成）"

    # If the markdown only has a single top-level title and all content is under it,
    # prefer the section structure for the deck title.
    if title == "需求描述":
        title = "智能派单需求说明"

    bg = Path(args.bg) if args.bg else None
    bg_title = Path(args.bg_title) if args.bg_title else None
    asset_dir = out_path.parent / ".ppt_assets"
    themes = ensure_builtin_themes(asset_dir)
    if args.list_themes:
        print("Built-in themes:")
        for t in themes:
            print(f"  - {t}: {asset_dir / 'themes' / t}")
        return 0

    build_ppt(
        sections,
        title=title,
        subtitle=subtitle,
        output_path=out_path,
        bg_path=bg,
        bg_title_path=bg_title,
        theme=args.theme,
    )
    print(f"OK: {out_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
